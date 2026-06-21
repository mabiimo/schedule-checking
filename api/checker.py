"""
Core logic untuk mencocokkan data Plan (Data1) vs Actual (Data2).

Aturan pencocokan:
- Kunci pencocokan: (tanggal, jam, brand)
- Jika kombinasi ditemukan di actual -> ambil status dari kolom `status` di actual,
  dinormalisasi jadi "Selesai" atau "Belum".
- Jika kombinasi plan TIDAK ditemukan sama sekali di actual -> status = "Tidak Ditemukan".
"""
from __future__ import annotations

import io
from datetime import datetime
from typing import Any

import pandas as pd

REQUIRED_PLAN_COLUMNS = {"tanggal", "start jam", "brand"}
REQUIRED_ACTUAL_COLUMNS = {"tanggal", "start jam", "brand", "status"}

STATUS_SELESAI = "Selesai"
STATUS_BELUM = "Belum"
STATUS_TIDAK_DITEMUKAN = "Tidak Ditemukan"

FILTER_SEMUA = "Semua"
VALID_STATUS_FILTERS = {FILTER_SEMUA, STATUS_SELESAI, STATUS_BELUM, STATUS_TIDAK_DITEMUKAN}


class DataValidationError(Exception):
    """Dilempar saat struktur file Excel yang di-upload tidak sesuai."""


def filter_by_status(result_df: pd.DataFrame, status_filter: str | None) -> pd.DataFrame:
    """Filter hasil cek berdasarkan status. None atau 'Semua' artinya tidak difilter."""
    if not status_filter or status_filter == FILTER_SEMUA:
        return result_df

    if status_filter not in VALID_STATUS_FILTERS:
        raise DataValidationError(
            f"Filter status '{status_filter}' tidak dikenali. "
            f"Gunakan salah satu dari: {', '.join(sorted(VALID_STATUS_FILTERS))}."
        )

    return result_df[result_df["status"] == status_filter].reset_index(drop=True)


def _normalize_columns(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()
    df.columns = [str(c).strip().lower() for c in df.columns]
    return df


def _normalize_tanggal(series: pd.Series) -> pd.Series:
    parsed = pd.to_datetime(series, errors="coerce")
    return parsed.dt.strftime("%Y-%m-%d")


def _normalize_text(series: pd.Series) -> pd.Series:
    return series.astype(str).str.strip()


def _normalize_status(value: str) -> str:
    value = (value or "").strip().lower()
    if value in {"selesai", "done", "ok", "complete", "completed"}:
        return STATUS_SELESAI
    return STATUS_BELUM


def load_and_validate(file_bytes: bytes, required_columns: set[str], label: str) -> pd.DataFrame:
    try:
        df = pd.read_excel(io.BytesIO(file_bytes))
    except Exception as exc:
        raise DataValidationError(f"Gagal membaca file {label}: file bukan Excel yang valid.") from exc

    df = _normalize_columns(df)
    missing = required_columns - set(df.columns)
    if missing:
        raise DataValidationError(
            f"File {label} tidak memiliki kolom wajib: {', '.join(sorted(missing))}. "
            f"Kolom yang ditemukan: {', '.join(df.columns)}."
        )

    if df.empty:
        raise DataValidationError(f"File {label} tidak memiliki baris data.")

    df["tanggal"] = _normalize_tanggal(df["tanggal"])
    if df["tanggal"].isna().any():
        raise DataValidationError(f"File {label} memiliki nilai tanggal yang tidak valid.")

    df["start jam"] = _normalize_text(df["start jam"])
    df["brand"] = _normalize_text(df["brand"])

    return df


def compare_plan_actual(plan_bytes: bytes, actual_bytes: bytes) -> pd.DataFrame:
    plan_df = load_and_validate(plan_bytes, REQUIRED_PLAN_COLUMNS, "Plan (Data1)")
    actual_df = load_and_validate(actual_bytes, REQUIRED_ACTUAL_COLUMNS, "Actual (Data2)")

    actual_df["status_norm"] = actual_df["status"].apply(_normalize_status)

    actual_lookup = actual_df.set_index(["tanggal", "start jam", "brand"])["status_norm"]
    # Jika ada duplikat kunci di actual, ambil yang terakhir.
    actual_lookup = actual_lookup[~actual_lookup.index.duplicated(keep="last")]

    results: list[dict[str, Any]] = []
    for _, row in plan_df.iterrows():
        key = (row["tanggal"], row["start jam"], row["brand"])
        if key in actual_lookup.index:
            status = actual_lookup.loc[key]
        else:
            status = STATUS_TIDAK_DITEMUKAN

        results.append(
            {
                "tanggal": row["tanggal"],
                "start jam": row["start jam"],
                "brand": row["brand"],
                "status": status,
            }
        )

    result_df = pd.DataFrame(results, columns=["tanggal", "start jam", "brand", "status"])
    result_df.sort_values(by=["brand", "tanggal", "start jam"], inplace=True, ignore_index=True)
    return result_df


def build_summary(result_df: pd.DataFrame) -> pd.DataFrame:
    """Ringkasan jumlah Selesai/Belum/Tidak Ditemukan per brand."""
    if result_df.empty:
        return pd.DataFrame(columns=["brand", STATUS_SELESAI, STATUS_BELUM, STATUS_TIDAK_DITEMUKAN, "total"])

    pivot = (
        result_df.groupby(["brand", "status"])
        .size()
        .unstack(fill_value=0)
        .reindex(columns=[STATUS_SELESAI, STATUS_BELUM, STATUS_TIDAK_DITEMUKAN], fill_value=0)
    )
    pivot["total"] = pivot.sum(axis=1)
    pivot = pivot.reset_index()
    return pivot


def export_to_excel(result_df: pd.DataFrame, summary_df: pd.DataFrame) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    wb = Workbook()

    ws = wb.active
    ws.title = "Hasil Cek"

    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill("solid", start_color="305496")
    selesai_fill = PatternFill("solid", start_color="C6EFCE")
    belum_fill = PatternFill("solid", start_color="FFEB9C")
    tidak_ditemukan_fill = PatternFill("solid", start_color="FFC7CE")

    headers = ["Tanggal", "Start Jam", "Brand", "Status"]
    ws.append(headers)
    for col_idx in range(1, len(headers) + 1):
        cell = ws.cell(row=1, column=col_idx)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center")

    status_fill_map = {
        STATUS_SELESAI: selesai_fill,
        STATUS_BELUM: belum_fill,
        STATUS_TIDAK_DITEMUKAN: tidak_ditemukan_fill,
    }

    for _, row in result_df.iterrows():
        ws.append([row["tanggal"], row["start jam"], row["brand"], row["status"]])
        status_cell = ws.cell(row=ws.max_row, column=4)
        fill = status_fill_map.get(row["status"])
        if fill:
            status_cell.fill = fill

    for col_idx, width in enumerate([14, 18, 16, 16], start=1):
        ws.column_dimensions[get_column_letter(col_idx)].width = width

    ws.freeze_panes = "A2"

    ws2 = wb.create_sheet("Ringkasan")
    summary_headers = list(summary_df.columns)
    ws2.append(summary_headers)
    for col_idx in range(1, len(summary_headers) + 1):
        cell = ws2.cell(row=1, column=col_idx)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center")

    for _, row in summary_df.iterrows():
        ws2.append(list(row))

    for col_idx in range(1, len(summary_headers) + 1):
        ws2.column_dimensions[get_column_letter(col_idx)].width = 16

    buffer = io.BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


PPTX_ROWS_PER_SLIDE = 10


def export_to_pptx(result_df: pd.DataFrame, summary_df: pd.DataFrame) -> bytes:
    """
    Buat file PPTX, dikelompokkan per brand:
    - 1 slide cover berisi nama brand (+ ringkasan singkat).
    - N slide tabel data, maksimal PPTX_ROWS_PER_SLIDE baris per slide.
    Tabel dibuat sederhana (tanpa pewarnaan status).
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    COLOR_BG = RGBColor(0x0B, 0x13, 0x20)
    COLOR_TEXT = RGBColor(0xF3, 0xF3, 0xEE)
    COLOR_MUTED = RGBColor(0x9A, 0xA7, 0xC2)
    COLOR_HEADER_FILL = RGBColor(0x30, 0x54, 0x96)
    COLOR_HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
    COLOR_ROW_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
    COLOR_ROW_ALT = RGBColor(0xF2, 0xF2, 0xF2)

    def add_cover_slide(brand: str, brand_summary_row: dict[str, Any] | None):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG

        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(11.3), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = brand
        run.font.size = Pt(54)
        run.font.bold = True
        run.font.color.rgb = COLOR_TEXT
        run.font.name = "Calibri"

        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.1), Inches(11.3), Inches(0.6))
        sub_tf = sub_box.text_frame
        sub_p = sub_tf.paragraphs[0]
        sub_p.alignment = PP_ALIGN.LEFT
        sub_run = sub_p.add_run()
        sub_run.text = "Cek Plan vs Actual"
        sub_run.font.size = Pt(16)
        sub_run.font.color.rgb = COLOR_MUTED
        sub_run.font.name = "Calibri"

        if brand_summary_row is not None:
            stats_text = (
                f"{brand_summary_row.get(STATUS_SELESAI, 0)} selesai   ·   "
                f"{brand_summary_row.get(STATUS_BELUM, 0)} belum   ·   "
                f"{brand_summary_row.get(STATUS_TIDAK_DITEMUKAN, 0)} tidak ditemukan   ·   "
                f"{brand_summary_row.get('total', 0)} total baris"
            )
            stats_box = slide.shapes.add_textbox(Inches(1), Inches(4.7), Inches(11.3), Inches(0.6))
            stats_tf = stats_box.text_frame
            stats_p = stats_tf.paragraphs[0]
            stats_p.alignment = PP_ALIGN.LEFT
            stats_run = stats_p.add_run()
            stats_run.text = stats_text
            stats_run.font.size = Pt(14)
            stats_run.font.color.rgb = COLOR_MUTED
            stats_run.font.name = "Calibri"

    def add_table_slide(brand: str, chunk_df: pd.DataFrame, page_no: int, total_pages: int):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        header_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(10), Inches(0.5))
        header_tf = header_box.text_frame
        header_p = header_tf.paragraphs[0]
        header_run = header_p.add_run()
        header_run.text = f"{brand}"
        header_run.font.size = Pt(20)
        header_run.font.bold = True
        header_run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
        header_run.font.name = "Calibri"

        page_box = slide.shapes.add_textbox(Inches(10.8), Inches(0.3), Inches(2), Inches(0.5))
        page_tf = page_box.text_frame
        page_p = page_tf.paragraphs[0]
        page_p.alignment = PP_ALIGN.RIGHT
        page_run = page_p.add_run()
        page_run.text = f"Halaman {page_no} / {total_pages}"
        page_run.font.size = Pt(12)
        page_run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
        page_run.font.name = "Calibri"

        n_rows = len(chunk_df) + 1
        n_cols = 4
        table_left = Inches(0.6)
        table_top = Inches(1.0)
        table_width = Inches(12.1)
        table_height = Inches(0.5 * n_rows)

        graphic_frame = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_width, table_height)
        table = graphic_frame.table

        col_widths = [Inches(2.6), Inches(3.0), Inches(2.6), Inches(3.9)]
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = width

        headers = ["Tanggal", "Start Jam", "Brand", "Status"]
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_HEADER_FILL
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            run = para.runs[0]
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = COLOR_HEADER_TEXT
            run.font.name = "Calibri"

        for row_offset, (_, row) in enumerate(chunk_df.iterrows(), start=1):
            values = [row["tanggal"], row["start jam"], row["brand"], row["status"]]
            for col_idx, value in enumerate(values):
                cell = table.cell(row_offset, col_idx)
                cell.text = str(value)
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_ROW_ALT if row_offset % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
                para = cell.text_frame.paragraphs[0]
                para.alignment = PP_ALIGN.LEFT
                run = para.runs[0]
                run.font.size = Pt(12)
                run.font.color.rgb = COLOR_ROW_TEXT
                run.font.name = "Calibri"

    if result_df.empty:
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
        p = box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Tidak ada data untuk ditampilkan."
        run.font.size = Pt(28)
        run.font.color.rgb = COLOR_TEXT
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    summary_lookup = {row["brand"]: row for row in summary_df.to_dict(orient="records")}

    for brand in sorted(result_df["brand"].unique()):
        brand_df = result_df[result_df["brand"] == brand].reset_index(drop=True)
        add_cover_slide(brand, summary_lookup.get(brand))

        total_pages = max(1, (len(brand_df) + PPTX_ROWS_PER_SLIDE - 1) // PPTX_ROWS_PER_SLIDE)
        for page_idx in range(total_pages):
            start = page_idx * PPTX_ROWS_PER_SLIDE
            end = start + PPTX_ROWS_PER_SLIDE
            chunk = brand_df.iloc[start:end]
            add_table_slide(brand, chunk, page_idx + 1, total_pages)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def generate_output_filename() -> str:
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"hasil_cek_plan_actual_{timestamp}.xlsx"


def generate_pptx_filename() -> str:
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"hasil_cek_plan_actual_{timestamp}.pptx"
